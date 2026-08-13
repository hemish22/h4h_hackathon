"""Generate the hackathon slide deck -> Hack4hackathon_presentation.pptx
Embeds the real figures from figures/ and the real result numbers.
Run: python scripts/make_ppt.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "figures")

INK = RGBColor(0x22, 0x28, 0x31)
SLATE = RGBColor(0x39, 0x3E, 0x46)
YELLOW = RGBColor(0xFF, 0xD3, 0x69)
MIST = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(4 * space)
        r = p.add_run(); r.text = text
        f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
        f.name = "Calibri"
    return tb


def _rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def content_header(slide, kicker, title):
    _bg(slide, MIST)
    _rect(slide, 0, 0, Inches(0.28), SH, YELLOW)
    _txt(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.4),
         [(kicker.upper(), 13, SLATE, True, False)])
    _txt(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.9),
         [(title, 30, INK, True, False)])


def bullets(slide, items, x=Inches(0.75), y=Inches(1.9), w=Inches(11.8), h=Inches(5.0), size=18):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        lead, rest = (it if isinstance(it, tuple) else (None, it))
        if lead:
            r = p.add_run(); r.text = lead + "  "; r.font.bold = True
            r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = "Calibri"
        r2 = p.add_run(); r2.text = rest
        r2.font.size = Pt(size); r2.font.color.rgb = SLATE; r2.font.name = "Calibri"


def title_slide():
    s = prs.slides.add_slide(BLANK); _bg(s, INK)
    _rect(s, Inches(0.9), Inches(2.55), Inches(1.7), Inches(0.12), YELLOW)
    _txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5),
         [("HACK4HACKATHON", 16, YELLOW, True, False)])
    _txt(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(2.0),
         [("Explainable Multimodal", 44, WHITE, True, False),
          ("Mental-Health Screening", 44, WHITE, True, False)])
    _txt(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.4),
         [("Face + voice + 18 behavioural/physiological signals → 4-class stress triage "
           "+ 3 severity scores, with per-modality confidence, disagreement detection, and a "
           "clinician-readable report card.", 17, MIST, False, False),
          ("A screening tool that flags people for professional assessment — not a diagnosis.",
           15, YELLOW, False, True)])


def section_slide(kicker, title, sub=None):
    s = prs.slides.add_slide(BLANK); _bg(s, INK)
    _rect(s, Inches(0.9), Inches(2.9), Inches(1.7), Inches(0.12), YELLOW)
    _txt(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
         [(kicker.upper(), 15, YELLOW, True, False)])
    _txt(s, Inches(0.85), Inches(3.2), Inches(11.6), Inches(1.6),
         [(title, 38, WHITE, True, False)])
    if sub:
        _txt(s, Inches(0.9), Inches(4.9), Inches(11.4), Inches(1.2),
             [(sub, 17, MIST, False, False)])


def image_slide(kicker, title, img, caption=None, img_w=Inches(10.6)):
    s = prs.slides.add_slide(BLANK); content_header(s, kicker, title)
    path = os.path.join(FIG, img)
    pic = s.shapes.add_picture(path, 0, 0, width=img_w)
    pic.left = int((SW - pic.width) / 2)
    pic.top = Inches(1.75)
    if caption:
        _txt(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.5),
             [(caption, 14, SLATE, False, True)], align=PP_ALIGN.CENTER)
    return s


def two_image_slide(kicker, title, imgs, caption=None):
    s = prs.slides.add_slide(BLANK); content_header(s, kicker, title)
    x = Inches(0.6)
    for img in imgs:
        pic = s.shapes.add_picture(os.path.join(FIG, img), x, Inches(2.0), width=Inches(6.0))
        x = Inches(6.8)
    if caption:
        _txt(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
             [(caption, 14, SLATE, False, True)], align=PP_ALIGN.CENTER)
    return s


def table_slide(kicker, title, header, rows, caption=None, highlight_row=None):
    s = prs.slides.add_slide(BLANK); content_header(s, kicker, title)
    nrows, ncols = len(rows) + 1, len(header)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.95),
                            Inches(12.1), Inches(0.4 * nrows)).table
    for j, htext in enumerate(header):
        c = gt.cell(0, j); c.text = htext
        c.fill.solid(); c.fill.fore_color.rgb = INK
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        hl = (highlight_row is not None and i - 1 == highlight_row)
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = YELLOW if hl else WHITE
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(12); r.font.color.rgb = INK; r.font.bold = (hl or j == 0)
    if caption:
        _txt(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.6),
             [(caption, 14, SLATE, False, True)], align=PP_ALIGN.CENTER)
    return s


# ============================================================ BUILD
title_slide()

s = prs.slides.add_slide(BLANK); content_header(s, "The problem", "Distress is measurable — but rarely measured")
bullets(s, [
    ("Today:", "psychiatric assessment leans on infrequent, self-reported questionnaires and interviews."),
    ("But:", "distress also shows in face, voice, sleep, social engagement, and physiology — continuously observable."),
    ("Our task (3 objectives):", "(1) classify stress into 4 levels, (2) estimate Depression/Anxiety/Stress scores, "
     "(3) explain and quantify each modality's contribution."),
    ("Framing:", "a screening / decision-support tool — flags people for a professional. Not a diagnosis."),
])

s = prs.slides.add_slide(BLANK); content_header(s, "Why multimodal", "Each channel alone is ambiguous")
bullets(s, [
    ("Flat voice", "could be a cold — or blunted affect."),
    ("A frown", "could be concentration — or distress."),
    ("The signal is agreement.", "When face, voice and physiology point the same way, confidence is high; "
     "when they disagree, that itself is information."),
    ("So we fuse the three — and, uniquely, we measure when they contradict.", ""),
])

section_slide("The core challenge", "Three datasets. Zero shared people.",
              "Faces (FER2013), voices (RAVDESS) and the 18-feature CSV describe different people. "
              "Yet the model needs one person with all three. Something had to be constructed.")

image_slide("Our differentiator", "Feature-matched alignment — joining a photo, a voice and 18 numbers",
            "alignment_diagram.png",
            "We DON'T staple a random face+voice to a CSV row — we match on measured features within the stress class.")

s = prs.slides.add_slide(BLANK); content_header(s, "Why it's principled", "The CSV told us what each participant should look & sound like")
bullets(s, [
    ("The CSV already contains", "audio-derived columns (MFCC_Mean, Pitch_Mean, Speech_Rate) and a face-derived one (Facial_Emotion_Variance)."),
    ("So we compute the SAME quantities", "from the raw wavs/images and match each CSV row to its nearest real voice & face."),
    ("Speech_Rate = 6 / clip-duration — exactly.", "Every RAVDESS sentence is 6 words (read off the spec, not guessed)."),
    ("We built BOTH", "a random-paired and a matched-paired dataset, and trained the identical model on each — to prove matching helps."),
])

s = prs.slides.add_slide(BLANK); content_header(s, "What the data told us", "Three honest findings (we report them, not hide them)")
bullets(s, [
    ("Tabular is near-noise:", "the 18 features barely correlate with the targets — LightGBM alone gets 37% (below guessing 'Healthy')."),
    ("Severe is rare:", "~3% of rows, only 19 in the test set → low Severe recall is a stated limitation, not a bug."),
    ("A contradiction in the committee's own labels:", "'angry' & 'disgust' map to OPPOSITE stress levels in the audio vs image tables — "
     "we flag every affected pair so it can't corrupt the analysis."),
])

image_slide("Architecture", "Three encoders → gated fusion → multi-task, uncertainty-aware",
            "architecture_diagram.png")

s = prs.slides.add_slide(BLANK); content_header(s, "Key techniques", "Small additions, large payoff")
bullets(s, [
    ("Ordinal-aware loss:", "predicting Healthy for a Severe case hurts more than predicting Moderate (the classes are an ordered scale)."),
    ("Per-participant gated fusion:", "'for this person the model leaned 41% on voice' — an explainability output, not just a mechanism."),
    ("Auxiliary heads:", "free single-modality results + the ingredients for decision-level fusion and concordance."),
    ("Modality dropout:", "the demo still works if a judge uploads only a photo."),
    ("MC-dropout:", "'Depression 21 ± 4', not a fake-precise 21.3 — shows what the model doesn't know."),
])

table_slide("Did we even need the pairing?", "Decision-level fusion — an honest baseline with NO pairing",
            ["Capability", "Decision-level (no pairing)", "Joint gated fusion"],
            [["4-class classification", "✓  (acc 0.652)", "✓  (acc 0.670)"],
             ["Depression/Anxiety/Stress regression", "✗  collapses to tabular", "✓"],
             ["Per-participant modality weights", "✗  global only", "✓"],
             ["Cross-modal representation learning", "✗", "✓"]],
            "Tuned weights: face 0.45 / voice 0.50 / tab 0.05 — the near-zero tab weight confirms tabular is noise. "
            "Objectives 2 & 3 REQUIRE the joint model → that's why the alignment engine exists.")

table_slide("Results", "The headline: matched pairing beats random — same model",
            ["Model", "Pairing", "Acc", "Macro F1", "ROC-AUC", "QWK", "Str MAE"],
            [["Tabular only (LightGBM)", "n/a", "0.370", "0.260", "0.493", "0.050", "10.11"],
             ["Gated fusion", "random", "0.553", "0.474", "0.775", "0.437", "9.28"],
             ["Gated fusion", "matched", "0.655", "0.571", "0.849", "0.564", "9.06"],
             ["Final (+ordinal +mod-dropout)", "matched", "0.670", "0.575", "0.864", "0.562", "8.95"]],
            "Random → matched = +10 accuracy points / +0.13 QWK under an IDENTICAL model → isolates the alignment engine's contribution.",
            highlight_row=3)

s = prs.slides.add_slide(BLANK); content_header(s, "Standout feature", "Modality disagreement → route to a human")
bullets(s, [
    ("We measure concordance:", "do face, voice and tabular agree? (masked affect: people hide distress in the face more than the voice.)"),
    ("It's a REAL reliability signal:", "on the clean test subset, accuracy is 0.64 when modalities agree — but only 0.23 when they strongly disagree."),
    ("So we route:", "high concordance → report; medium → caveat; low → FLAG FOR HUMAN REVIEW instead of a confident label."),
    ("Calibration:", "reliability diagram + ECE = 0.198 reported — we show how trustworthy the confidences are."),
])

two_image_slide("Explainability", "Where the model looked + what drives it",
                ["shap_beeswarm.png", "gradcam_faces.png"],
                "SHAP: HRV is the top tabular driver (matches stress physiology).  Grad-CAM: the model attends to eyes / nose / mouth.")

s = prs.slides.add_slide(BLANK); content_header(s, "The prototype", "Not just a predictor — a trust layer")
bullets(s, [
    ("Shows the actual inputs it joined:", "each one-click demo displays the real FACE image and a playable AUDIO clip that were paired — our alignment idea, made visible."),
    ("Per-participant gate weights, MC-dropout uncertainty bands, live Grad-CAM on your own face & voice.", ""),
    ("Counterfactuals:", "'what would change this' — the actionable movers."),
    ("Graceful degradation + webcam/mic capture + conflict-aware honesty messaging.", ""),
    ("'About the model' panel:", "ablation table, calibration, concordance evidence — rigour on demand."),
], size=17)

section_slide("Live demo", "Let's screen four participants",
              "Clear Healthy · Severe · a deliberately low-concordance case (caught and routed to review) · a label-conflict case (honesty behaviour).")

s = prs.slides.add_slide(BLANK); content_header(s, "Limitations", "We name them before you do")
bullets(s, [
    ("Synthetic pairing", "— constructed participants, not a real cohort; conclusions rest on RELATIVE comparisons."),
    ("Acted emotions", "(RAVDESS/FER) → live webcam/mic is out-of-distribution."),
    ("Skewed CSV", "— Severe ≈ 3% (19 test participants); we quote effective sample size."),
    ("Only 1 of 4 face-video columns recoverable", "from static stills (no blink/head-motion)."),
    ("Mapping conflict", "in the committee labels; ECE 0.198 = mild overconfidence. Not clinically validated.")
])

s = prs.slides.add_slide(BLANK); content_header(s, "Next steps & takeaway", "Where this goes")
bullets(s, [
    ("Real longitudinal cohort", "with shared participants → removes the synthetic-pairing caveat."),
    ("Temporal modelling", "→ recovers blink rate & head motion; calibration/temperature scaling to cut ECE."),
    ("On-device, privacy-preserving inference.", ""),
])
_txt(s, Inches(0.75), Inches(5.4), Inches(11.8), Inches(1.6),
     [("In one line:", 18, INK, True, False),
      ("Everyone stapled three encoders together — we aligned participants in measured-feature space and proved "
       "it helped (+10 pts), found the contradiction in the committee's own labels, treated the classes as the "
       "ordered scale they are, reported what the model doesn't know, and shipped something a clinician could read.",
       17, SLATE, False, True)])

out = os.path.join(ROOT, "Hack4hackathon_presentation.pptx")
prs.save(out)
print("wrote", out, "-", len(prs.slides._sldIdLst), "slides")
